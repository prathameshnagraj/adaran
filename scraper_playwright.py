# scraper_playwright_enhanced.py

import os
import json
import asyncio
import requests
import pandas as pd
from pathlib import Path
from playwright.async_api import async_playwright
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

# Config
INPUT_EXCEL = "all_jsom_urls_combined.xlsx"
PAGE_SAVE_DIR = "playwright_scraped/pages"
DOC_SAVE_DIR = "playwright_scraped/documents"
Path(PAGE_SAVE_DIR).mkdir(parents=True, exist_ok=True)
Path(DOC_SAVE_DIR).mkdir(parents=True, exist_ok=True)

HEADERS = {"User-Agent": "Mozilla/5.0"}

# Document Extractors
def extract_pdf_text(filepath):
    try:
        reader = PdfReader(filepath)
        return "\n".join([p.extract_text() for p in reader.pages if p.extract_text()])
    except Exception as e:
        print(f"⚠️ PDF extraction failed: {filepath}, {e}")
        return ""

def extract_docx_text(filepath):
    try:
        doc = DocxDocument(filepath)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        print(f"⚠️ DOCX extraction failed: {filepath}, {e}")
        return ""

def extract_pptx_text(filepath):
    try:
        prs = Presentation(filepath)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
    except Exception as e:
        print(f"⚠️ PPTX extraction failed: {filepath}, {e}")
        return ""

# Document Downloader
def download_and_extract(doc_url):
    try:
        response = requests.get(doc_url, headers=HEADERS, timeout=20)
        response.raise_for_status()

        file_name = os.path.basename(doc_url.split("?")[0])  # Clean up URL params
        save_path = os.path.join(DOC_SAVE_DIR, file_name)

        with open(save_path, "wb") as f:
            f.write(response.content)

        if file_name.endswith(".pdf"):
            text = extract_pdf_text(save_path)
        elif file_name.endswith(".docx"):
            text = extract_docx_text(save_path)
        elif file_name.endswith(".pptx"):
            text = extract_pptx_text(save_path)
        else:
            text = ""

        return {
            "url": doc_url,
            "file_name": file_name,
            "extracted_text": text.strip()
        }

    except Exception as e:
        print(f"❌ Failed to download or extract {doc_url}: {e}")
        return None

# Playwright Scraping
async def scrape_page(playwright, url):
    browser = await playwright.chromium.launch(headless=True)
    page = await browser.new_page()

    try:
        await page.goto(url, wait_until="networkidle", timeout=40000)
        text = await page.evaluate("document.body.innerText")

        # Find all <a> links
        links = await page.locator('a').evaluate_all("els => els.map(el => el.href)")
        document_links = [link for link in links if link.lower().endswith(('.pdf', '.docx', '.pptx'))]

        return {
            "page_data": {
                "url": url,
                "extracted_text": text.strip()
            },
            "document_links": document_links
        }

    except Exception as e:
        print(f"❌ Failed to scrape {url}: {e}")
        return None

    finally:
        await browser.close()

async def scrape_all(urls):
    page_texts = []
    documents_texts = []

    async with async_playwright() as playwright:
        for idx, url in enumerate(urls):
            print(f"🔍 Scraping {idx+1}/{len(urls)}: {url}")
            result = await scrape_page(playwright, url)

            if not result:
                continue

            # Save page text
            page_texts.append(result["page_data"])

            # Download and parse documents
            for doc_url in result["document_links"]:
                doc_result = download_and_extract(doc_url)
                if doc_result:
                    documents_texts.append(doc_result)

    # Save JSONs
    with open(os.path.join(PAGE_SAVE_DIR, "scraped_pages.json"), "w", encoding="utf-8") as f:
        json.dump(page_texts, f, indent=2, ensure_ascii=False)

    with open(os.path.join(DOC_SAVE_DIR, "scraped_documents.json"), "w", encoding="utf-8") as f:
        json.dump(documents_texts, f, indent=2, ensure_ascii=False)

    print(f"✅ Scraping completed: {len(page_texts)} pages and {len(documents_texts)} documents extracted.")

# Entry Point
if __name__ == "__main__":
    df = pd.read_excel(INPUT_EXCEL)
    urls = df['URL'].dropna().tolist()

    asyncio.run(scrape_all(urls))
